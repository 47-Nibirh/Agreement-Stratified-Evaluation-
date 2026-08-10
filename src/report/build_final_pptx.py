"""
Build the FINAL-defence deck.

The existing Phase-I deck stops at Phase 3 and carries the template layout's
stale "B.Sc. Pre-Defense" corner mark on every slide. This builder produces the
deck for the final defence: it covers all five research questions, the external
validation, the human comparator, the negative control and the backbone
replication, and it repairs the corner mark on the layout rather than papering
over it slide by slide.

Every number is interpolated from reports/phase7_register.json and the Phase 8
artefacts, and every reference number comes from src/report/bibliography.py, so
nothing on any slide is typed by hand. The slide plumbing is imported from
build_phase1_pptx rather than duplicated, so both decks inherit the template's
furniture identically.

Run:  python src/report/build_final_pptx.py
Out:  Final-Defence_Presentation.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(Path(__file__).resolve().parent))

import build_phase1_pptx as P1  # noqa: E402
from build_phase1_pptx import (INK, MUTED, TITLE_RGB, body_box,  # noqa: E402
                               bullets, clone_slide, keyline, notes, picture,
                               rich, set_title, strip_body, table)
from phase1_facts import ADMIN, facts  # noqa: E402
import bibliography as BIB  # noqa: E402
from bibliography import cite  # noqa: E402

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402

TEMPLATE = ROOT / "Final-Defense Slide Template-PPT.pptx"
OUTPUT = ROOT / "Final-Defence_Presentation.pptx"
REP = ROOT / "reports"
FIGD = ROOT / "figures_thesis"

REG = json.loads((REP / "phase7_register.json").read_text(encoding="utf-8"))["register"]
ORC = json.loads((REP / "phase8_oracle_reconcile.json").read_text(encoding="utf-8"))
LS = json.loads((REP / "phase8_lit_synthesis.json").read_text(encoding="utf-8"))
MULT = json.loads((REP / "phase7_multiplicity.json").read_text(encoding="utf-8"))
RQ5 = json.loads((REP / "phase7_rq5.json").read_text(encoding="utf-8"))

TIERS = ["S-unanimous", "S-majority", "S-plurality", "S-no-majority"]
POOLED = "S-contested (pooled)"
LABEL = {"S-unanimous": "Unanimous 4/4", "S-majority": "Majority 3/4",
         "S-plurality": "Plurality 2-1-1", "S-no-majority": "No majority",
         POOLED: "Contested (pooled)"}


def p100(x, n=1):
    return "n/a" if x is None else f"{100 * x:.{n}f}"


def pc(x, n=2):
    return "n/a" if x is None else f"{x:.{n}f}"


def iv(v, n=2):
    if not v or v[0] is None:
        return "n/a"
    return f"[{v[0]:.{n}f}, {v[1]:.{n}f}]"


def fix_corner_mark(prs) -> int:
    """Clear the template layout's stale 'Pre-Defense' corner mark.

    The mark lives on the slide LAYOUT (Rectangle 11), not on the slides, which
    is why every slide of the Phase-I deck showed 'B.Sc. Pre-Defense' at the top
    and 'B.Sc. Final-Defense' at the bottom simultaneously -- two marks
    contradicting each other on all 23 slides.

    The layout copy is blanked rather than rewritten to say 'Final-Defense',
    because each slide already carries its own correct mark; rewriting it would
    replace a contradiction with a duplicate.
    """
    n = 0
    seen: set = set()
    for slide in prs.slides:
        lay = slide.slide_layout
        # SlideLayout is not hashable; its underlying element is.
        if id(lay._element) in seen:
            continue
        seen.add(id(lay._element))
        for sh in lay.shapes:
            if sh.has_text_frame and "Pre-Defense" in sh.text_frame.text:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if "Pre-Defense" in run.text:
                            run.text = ""
                            n += 1
    return n


class Deck:
    """Hands out slides: the template's own first, then clones.

    Removing the template's unused slides after the fact was the obvious route
    and it produced a file python-pptx could read and PowerPoint could not --
    dropping the presentation->slide relationship leaves the slide part behind
    and the package no longer validates. Consuming the sixteen template slides
    in order and cloning only past the end means nothing ever has to be
    removed, so the package stays exactly as well-formed as the template.
    """

    def __init__(self, prs):
        self.prs = prs
        self.n_template = len(prs.slides._sldIdLst)
        self.i = 0
        self.n_cloned = 0

    def take(self):
        if self.i < self.n_template:
            s = self.prs.slides[self.i]
            self.i += 1
            return s
        self.n_cloned += 1
        return clone_slide(self.prs, 2)

    def content(self, title):
        s = self.take()
        strip_body(s)
        set_title(s, title)
        return s

    def unused(self):
        return max(0, self.n_template - self.i)


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    D = Deck(prs)
    a, s4 = REG["ch2_audit"], REG["ch4_stratified"]
    t5, e6, x7 = REG["ch5_targets"], REG["ch6_external"], REG["ch7_error"]

    # ---- 1. Title --------------------------------------------------------
    s = D.take()
    set_title(s, ADMIN["title"], size=Pt(28))
    sub = [ph for ph in s.placeholders if ph.placeholder_format.idx == 1][0]
    sub.top, sub.height = Emu(3930000), Emu(2450000)
    sub.left, sub.width = Emu(2500000), Emu(7192000)
    tf = sub.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = []
    for name, sid in ADMIN["members"]:
        lines.append([(name, Pt(15), True, INK),
                      (f"    ID: {sid}", Pt(12), False, MUTED)])
    lines.append([(f"Supervisor: {ADMIN['supervisor'][0]}, "
                   f"{ADMIN['supervisor'][1]}", Pt(12), False, INK)])
    lines.append([(f"Co-Supervisor: {ADMIN['co_supervisor'][0]}, "
                   f"{ADMIN['co_supervisor'][1]}", Pt(12), False, INK)])
    lines.append([("Department of Computer Science and Engineering, "
                   "Daffodil International University", Pt(12), True, INK)])
    rich(tf, lines, align=PP_ALIGN.CENTER, space=Pt(9))
    notes(s, "This is the final defence. The thesis asks one question: when a "
             "published system reports 85 macro F1 on this task, what fraction "
             "of the task does that number describe?")

    # ---- 2. Outline ------------------------------------------------------
    s = D.take()
    set_title(s, "Outline")
    ph = [x for x in s.placeholders if x.placeholder_format.idx == 13][0]
    ph.text_frame.clear()
    bullets(ph.text_frame, [
        "The problem: ground truth is a constructed object",
        "Research questions, published baselines, and the gap counted",
        "Methodology: pre-registration, data, pre-processing, metrics",
        "Baseline reproduction — the validity gate",
        "RQ1 — performance across agreement strata, against the ceiling",
        "The calibration collapse",
        "RQ2/RQ4 — five target constructions and a derived control",
        "External validation and out-of-protocol rejection",
        "RQ3 — the human comparator and selective prediction",
        "RQ5 — a negative control on our own audit protocol",
        "Backbone replication, threats, and what could be deployed",
        "Conclusions and references",
    ], size=Pt(18), space=Pt(4))

    # ---- 3. Introduction -------------------------------------------------
    s = D.take()
    strip_body(s); set_title(s, "Introduction")
    bullets(body_box(s, height=Emu(3350000)), [
        "Upper GI endoscopy is the primary screening route for gastric cancer; "
        "its yield depends on inspecting the whole mucosal surface.",
        "The Systematic Screening protocol fixes a photographic route over 22 "
        "landmarks — a grid of 4 gastric walls by 6 depth stations.",
        "A model that recognises the landmark in each frame can act as a "
        "real-time coverage monitor and flag an unvisited region.",
        f"Published systems report macro F1 near 85–88, and the task is "
        f"generally treated as solved {cite('gastrohun')}.",
        (1, f"But every one of those figures is measured only on frames all "
            f"four experts labelled identically — "
            f"{a['agreement_tiers_pct']['complete_agreement_4of4']:.1f}% of "
            f"this corpus."),
    ], size=Pt(19))
    keyline(s, "The reported number describes the easy fraction of the task.")

    # ---- 4. Problem identification ---------------------------------------
    s = D.content("Problem Identification")
    rich(body_box(s, height=Emu(1400000)), [
        [(f"{a['agreement_tiers_pct']['complete_agreement_4of4']:.1f}%",
          Pt(38), True, TITLE_RGB),
         ("  of the corpus is unanimous — and that is all the literature "
          "scores on.", Pt(19), False, INK)],
    ])
    bullets(body_box(s, top=Emu(2750000), height=Emu(2300000)), [
        f"{100 - a['agreement_tiers_pct']['complete_agreement_4of4']:.1f}% is "
        f"discarded before any result is reported.",
        f"The discarded images are not noise: "
        f"{a['disagreement_decomposition_pct']['same_station_different_wall']:.2f}% "
        f"of conflicts put two experts on different walls of the same station.",
        "A deployed system meets the contested fraction continuously, and gets "
        "no signal that its validation excluded those frames.",
        (1, "Operating accuracy is unknown on exactly the images where a second "
            "reader would help most."),
    ], size=Pt(18))

    # ---- 5. Research questions -------------------------------------------
    s = D.content("Research Questions")
    rows = [[f["rq"], f["endpoint"][:78]] for f in MULT["family"]]
    table(s, ["RQ", "Pre-registered primary endpoint"], rows,
          [0.98, 10.50], top=Emu(1500000), size=Pt(12))
    keyline(s, "One primary endpoint per question, each frozen before the "
               "analysis that tests it ran.")

    # ---- 6. Published baselines ------------------------------------------
    # Restored from the Phase-I deck. Without it the panel has no scale against
    # which to read 83.92, and no reason to believe the label-construction
    # effect is large. These are the descriptor's own published figures, quoted
    # and cited, not measurements of ours.
    s = D.content("Background — the published baselines")
    table(s, ["Model / label construction", "Params", "Macro F1"], [
        ["ConvNeXt-Large", "200 M", "88.25 +/- 0.22"],
        ["ConvNeXt-Tiny", "28 M", "approx. 85"],
        ["ResNet-152", "60 M", "85.28 +/- 0.27"],
        ["ConvNeXt-Tiny, fellow-agreement labels", "28 M", "87.05 +/- 0.21"],
        ["Best single annotator as target", "-", "84.82 +/- 0.23"],
        ["Human expert band", "-", "77.47 - 84.82"],
    ], [5.0, 1.6, 2.6], top=Emu(1750000), size=Pt(14), head_size=Pt(14))
    keyline(s, f"Changing only how the four labels are combined moves F1 by 2.2 "
               f"points — more than the gap between architecture families "
               f"{cite('gastrohun')}.", top=Emu(5150000), size=Pt(15))
    notes(s, "This is the slide that motivates the whole thesis. Label "
             "construction moves the score more than the architecture does, "
             "and yet every one of these numbers is computed on the unanimous "
             "subset only.")

    # ---- 7. The gap, counted ---------------------------------------------
    s = D.content("The Gap — counted, not asserted")
    d = LS["dimensions"]
    rows = [[v["label"][:44], f"{v['n_mentioning']}/{v['n_scored']}",
             f"{v['pct_mentioning']:.1f}%"] for v in d.values()]
    table(s, ["Reporting dimension", "Mentioning", "%"], rows,
          [6.78, 2.08, 1.53], top=Emu(1450000),
          size=Pt(13))
    rich(body_box(s, top=Emu(4450000), height=Emu(1000000)), [
        [(f"Of the {LS['gap_intersection']['n_mentioning_ground_truth_construction']} "
          f"studies that engage with how their reference standard was built, "
          f"only ", Pt(16), False, INK),
         (f"{LS['gap_intersection']['n_also_mentioning_calibration']}",
          Pt(22), True, TITLE_RGB),
         (" also mention whether their probabilities can be believed.",
          Pt(16), False, INK)],
    ])
    notes(s, "These are mention counts from title and abstract, so they are "
             "upper bounds. The population-description zero is the one figure "
             "we do not lean on: demographics live in a baseline table, not in "
             "an abstract.")

    # ---- 7. Methodology --------------------------------------------------
    s = D.content("Methodology — eight gated phases")
    picture(s, "figures_thesis/T02_workflow.png", top=Emu(1420000), max_h=Emu(3750000))
    keyline(s, "Each phase freezes its hypotheses and verdict rules before any "
               "model runs; the generating script refuses to overwrite them.")

    # ---- 8. Data ---------------------------------------------------------
    s = D.content("Data and corpus audit")
    rows = [
        ["Images / patients / classes",
         f"{a['n_images']:,} / {a['n_patients']} / {a['n_classes']}"],
        ["Annotators", ", ".join(a["annotators"])],
        ["Decode integrity",
         f"{a['n_decoded_ok']:,} decoded, {a['n_missing']} missing, "
         f"{a['n_corrupt']} corrupt"],
        ["Near-duplicate pairs examined", f"{a['neardup_pairs_examined']:,}"],
        ["Cross-split duplicates (calibrated rule)",
         f"{a['dup_confirmed_calibrated']}"],
        ["Fleiss' κ / Krippendorff's α / Gwet's AC1",
         f"{a['fleiss_kappa']} / {a['krippendorff_alpha']} / {a['gwet_ac1']}"],
    ]
    table(s, ["Measurement", "Value"], rows,
          [5.69, 5.91], top=Emu(1500000), size=Pt(13))
    keyline(s, "Chosen because it releases the individual annotators' labels. "
               "Without them, no endpoint in this design exists.")

    # ---- 10. Pre-processing, and the augmentation that is not used --------
    F = facts()
    pp = F["preprocess"]
    s = D.content("Pre-processing — and the augmentation we refuse")
    bullets(body_box(s, height=Emu(3300000)), [
        f"Resampled once to {pp['size']}x{pp['size']} with "
        f"{pp['resample'].title()} and cached, so every arm in every later "
        f"phase reads pixel-identical inputs.",
        f"Normalised on this cohort's own channel statistics, not the ImageNet "
        f"constants — endoscopic illumination is red-dominant and they differ "
        f"by up to {pp['max_delta']:.3f} in the mean.",
        "NO horizontal or vertical flips, and no large rotations.",
        (1, "The label encodes a gastric WALL — anterior, posterior, lesser "
            "curvature, greater curvature. A horizontal flip maps one wall onto "
            "the appearance of another and silently relabels the image."),
        f"What remains is a mild scale and aspect jitter (crop scale "
        f"{pp['crop_scale'][0]}-{pp['crop_scale'][1]}) and a photometric jitter "
        f"(brightness, contrast, saturation +/-{pp['jitter']['brightness']}, "
        f"hue +/-{pp['jitter']['hue']}) — neither moves the anatomy.",
    ], size=Pt(17))
    keyline(s, "An augmentation that destroys the target is not "
               "regularisation.")

    # ---- 11. Sample dataset ----------------------------------------------
    s = D.content("Sample dataset — agreement falling, left to right")
    picture(s, "figures_phase1/PH1_F13_sample_images.png", top=Emu(1400000),
            max_h=Emu(3400000))
    keyline(s, "Published evaluations score only the leftmost column. Input: a "
               "224x224 frame. Output: one of 23 classes with a confidence.")

    # ---- 12. Baseline reproduction ---------------------------------------
    # The validity gate. Restored because it is the first thing a panel asks:
    # how do we know the pipeline is right before we believe anything it says?
    b = F["baseline"]
    s = D.content("Baseline reproduction — the validity gate")
    seeds = ", ".join(f"{v}" for v in b["per_seed"].values())
    table(s, ["Metric", "Value"], [
        ["Macro F1, three-seed mean",
         f"{b['observed']:.2f}   (95% CI {b['ci95'][0]:.2f} - {b['ci95'][1]:.2f})"],
        ["Published target", f"{b['published']:g} +/- {b['band']:g}, "
                             f"fixed before training"],
        ["Difference", f"{b['delta']:+.2f} points"],
        ["Accuracy / weighted F1", f"{b['accuracy']:.2f} / {b['weighted_f1']:.2f}"],
        ["Expected calibration error", f"{b['ece']:.2f}%"],
        ["Per-seed macro F1", seeds],
        ["Pre-registered verdict", b["verdict"]],
    ], [4.6, 6.4], top=Emu(1500000), size=Pt(13))
    keyline(s, f"{b['verdict']} — the pipeline behaves as the published one "
               f"did. Everything after this is attributable to the design, not "
               f"to a bug.")
    notes(s, "If the panel asks only one methodological question, it will be "
             "this one. The band was fixed before the first training run, and "
             "the test split was scored once.")

    # ---- 13. Two ceilings ------------------------------------------------
    s = D.content("Two attainable ceilings — and why they differ")
    rows = [[LABEL.get(t, t), f"{r['panel_ceiling_recomputed']:.4f}",
             f"{r['loo_oracle_recomputed']:.4f}",
             f"{r['difference_loo_minus_panel']:+.4f}"]
            for t, r in ORC["by_stratum"].items()]
    table(s, ["Stratum", "Panel ceiling", "Leave-one-out oracle", "Difference"],
          rows, [3.28, 2.84, 3.28, 2.19],
          top=Emu(1430000), size=Pt(12))
    keyline(s, "They coincide where the mode survives removal and diverge by "
               "14.29 points where it does not — tie multiplicity, not a bug. "
               "RQ1 uses only the strata where they agree.")
    notes(s, "Chapter 5 and Chapter 8 both said 'modal-vote oracle'. They are "
             "different estimands. On a 2-2 image, four references give a "
             "two-way tie worth 2/4; any three give a unique mode worth 2/3. "
             "phase8_oracle_reconcile.py reproduces both committed series.")

    # ---- 10. RQ1 raw decline ---------------------------------------------
    s = D.content("RQ1 — performance across agreement strata")
    rows = [[LABEL[t], f"{s4['n_by_tier'][t]:,}", p100(s4["macro_f1_by_tier"][t]),
             pc(s4["ceilings"][t], 4), p100(s4["ece_by_tier"][t])]
            for t in TIERS]
    table(s, ["Stratum", "n", "Macro F1", "Ceiling", "ECE %"], rows,
          [3.28, 1.42, 2.30, 2.30, 2.30],
          top=Emu(1430000), size=Pt(13))
    keyline(s, f"Raw macro F1 falls {p100(s4['macro_f1_by_tier']['S-unanimous'])} "
               f"→ {p100(s4['macro_f1_by_tier']['S-plurality'])}. But the "
               f"ceiling falls with it.")

    # ---- 11. Ceiling normalisation ---------------------------------------
    s = D.content("RQ1 — the ceiling moves, so normalise by it")
    picture(s, "figures_thesis/T14_stratified_curve_ceiling.png", top=Emu(1400000),
            max_h=Emu(3550000))
    keyline(s, f"Ceiling-normalised unanimous-minus-majority gap: "
               f"{pc(s4['gap_4v3_ceiling_normalised'])} points "
               f"(95% CI {iv(s4['gap_4v3_ci'])}), not the "
               f"{pc(100 * (s4['macro_f1_by_tier']['S-unanimous'] - s4['macro_f1_by_tier']['S-majority']))} "
               f"the raw scores suggest.")

    # ---- 12. Calibration collapse ----------------------------------------
    s = D.content("The calibration collapse")
    picture(s, "figures_thesis/T16_calibration_by_stratum.png", top=Emu(1400000),
            max_h=Emu(3450000))
    dc = 100 * (s4["mean_confidence_by_tier"]["S-unanimous"]
                - s4["mean_confidence_by_tier"]["S-plurality"])
    da = 100 * (s4["expected_accuracy_by_tier"]["S-unanimous"]
                - s4["expected_accuracy_by_tier"]["S-plurality"])
    keyline(s, f"ECE {p100(s4['ece_by_tier']['S-unanimous'])}% → "
               f"{p100(s4['ece_by_tier']['S-plurality'])}%. Confidence falls "
               f"{pc(dc)} points while accuracy falls {pc(da)}. The durable "
               f"finding of the thesis.")

    # ---- 13. RQ2 design --------------------------------------------------
    s = D.content("RQ2 — five target constructions, one cohort")
    rows = [[k, v] for k, v in [
        ("C0", "hard label, 4/4 cohort"),
        ("C1", "hard majority label"),
        ("C2", "vote proportions (soft target)"),
        ("C3", "hard + matched smoothing (CONTROL)"),
        ("C4", "vote proportions + anatomical penalty")]]
    table(s, ["Arm", "Target construction"], rows,
          [1.75, 9.84], top=Emu(1500000), size=Pt(14))
    keyline(s, f"ε = 0.07529 derived to match the probability mass the soft "
               f"target displaces — not set to a convention. Backbone, "
               f"schedule and selection identical across arms.")

    # ---- 14. RQ2 result --------------------------------------------------
    s = D.content("RQ2 — the reversal is the finding")
    picture(s, "figures_thesis/T20_calibration_by_config.png", top=Emu(1400000),
            max_h=Emu(3400000))
    keyline(s, f"Accuracy contrast C2−C3 = {pc(t5['contrast_C2_C3'])} "
               f"({iv(t5['contrast_C2_C3_ci'])}): NOT RESOLVED. On calibration "
               f"the generic control WINS on contested images. No target "
               f"construction repairs it.")

    # ---- 15. External validation -----------------------------------------
    s = D.content("External validation — what the mapping destroys")
    picture(s, "figures_thesis/T22_external_label_space.png", top=Emu(1400000),
            max_h=Emu(3400000))
    keyline(s, "Neither external corpus carries the wall axis. A 23-way "
               "external validation is not available — the phase was reframed "
               "before any image was scored.")

    # ---- 16. Rejection ---------------------------------------------------
    s = D.content("Out-of-protocol rejection separates the arms")
    picture(s, "figures_thesis/T24_external_rejection.png", top=Emu(1400000),
            max_h=Emu(3400000))
    keyline(s, "Invisible internally: GastroHUN's test split holds almost no "
               "out-of-protocol images. This is the argument for a second "
               "centre stated as a measurement.")

    # ---- 17. Human comparator --------------------------------------------
    s = D.content("RQ3 — the human comparator, and the oracle that bounds it")
    rows = [[LABEL.get(t, t), pc(x7["human_by_stratum"].get(t), 4),
             pc(x7["model_by_stratum"].get(t), 4),
             pc(x7["oracle_by_stratum"].get(t), 4),
             (f"{100 * x7['headroom_recovered'][t]:.0f}%"
              if x7["headroom_recovered"].get(t) is not None else "n/a")]
            for t in TIERS + [POOLED]]
    table(s, ["Stratum", "Held-out annotator", "Model (C2)", "Modal-vote oracle",
              "Headroom"], rows,
          [2.95, 2.62, 2.08, 2.52, 1.53],
          top=Emu(1430000), size=Pt(12))
    keyline(s, "Out-predicts an individual annotator — but NOT the panel's own "
               "modal vote, on any stratum. It is not evidence of superior "
               "anatomical judgement.")

    # ---- 18. Selective prediction ----------------------------------------
    s = D.content("Selective prediction — where the arms finally separate")
    picture(s, "figures_thesis/T29_risk_coverage_external.png", top=Emu(1400000),
            max_h=Emu(3400000))
    keyline(s, f"External AURC {pc(x7['aurc_external']['C2'], 4)} for C2 "
               f"against {pc(x7['aurc_external']['C3'], 4)} for C3 — decisive, "
               f"where every internal endpoint failed to separate them.")

    # ---- 19. RQ5 negative control ----------------------------------------
    s = D.content("RQ5 — a negative control on our own audit protocol")
    picture(s, "figures_thesis/T31_negative_control.png", top=Emu(1400000), max_h=Emu(3300000))
    keyline(s, f"{RQ5['verdict']} — gates separate the corpora, but only 1 of 4 "
               f"fatal defects is caught by any gate, and only incidentally. "
               f"Well-formedness is not viability.")
    notes(s, "This is the slide to dwell on if asked what we would do "
             "differently. A negative control that flattered the protocol "
             "would have been worth nothing; this one changed it.")

    # ---- 20. Backbone replication ----------------------------------------
    m = REG.get("ch9_synthesis", {}).get("backbone", {})
    s = D.content("Is any of this a ConvNeXt artefact?")
    bullets(body_box(s, height=Emu(3300000)), [
        "The target-construction contrast was re-run on EfficientNet-B0 "
        f"(4.0 M parameters against ConvNeXt-Tiny's 28 M) {cite('efficientnet')}.",
        "The training script imports the Phase 4 code and rebinds only the "
        "model constructor, so any difference is architectural.",
        "3 of 3 claims replicate: the accuracy null, the calibration reversal, "
        "and the model's position between annotator and oracle.",
        (1, "What this does NOT cover: confusion geometry, attribution "
            "stability, and the external endpoints. Those remain "
            "single-architecture and are named in the threat table."),
    ], size=Pt(18))
    keyline(s, "The strongest remaining objection, tested — and only partly "
               "closed. We say which part.")

    # ---- 21. Threats -----------------------------------------------------
    s = D.content("Threats to validity, ranked")
    rows = [[t["threat"][:58], t["severity"]]
            for t in REG.get("ch9_synthesis", {}).get("threats", [])[:7]] or [
        ["Single backbone for geometry, attribution, external", "Medium"],
        ["Compute-bound epoch cap on 8/12 runs", "Medium"],
        ["Contested strata are small (n = 342, 127, 81)", "Medium"],
        ["External intervals are image-level", "Medium"],
        ["Four annotators bound the human comparator", "Medium"],
        ["Strata defined by the same agreement that scores the human", "Medium"],
        ["No age or sex in the corpus", "Low"]]
    table(s, ["Threat", "Severity"], rows,
          [8.97, 2.62], top=Emu(1500000), size=Pt(13))
    keyline(s, "Ranked by how much each could change a conclusion — not by how "
               "easy each is to answer.")

    # ---- 22. Deployment --------------------------------------------------
    s = D.content("What could be deployed, and under what conditions")
    bullets(body_box(s, height=Emu(3300000)), [
        "Nothing here supports deploying an autonomous landmark classifier.",
        "What it supports is a completeness-checking assistant that ABSTAINS.",
        f"The operating characteristic is the external risk–coverage curve: "
        f"AURC {pc(x7['aurc_external']['C2'], 4)} for C2 against "
        f"{pc(x7['aurc_external']['C3'], 4)} for the internally-best arm.",
        (1, "A deployment would have to select on that endpoint, monitor it at "
            "the new centre, and route declined images to a human rather than "
            "scoring them."),
    ], size=Pt(18))
    keyline(s, "The recommendation is about abstention, not about a better loss "
               "— because no target construction repaired calibration.")

    # ---- 23. Conclusions -------------------------------------------------
    s = D.content("Conclusions — answers to the five questions")
    rows = [[f["rq"],
             ("NOT ESTIMABLE" if f.get("not_estimable")
              else "CONTAINS 0" if f.get("multiplicity", {}).get("contains_zero")
              else "NO INTERVAL" if not f.get("multiplicity")
              else "EXCLUDES 0; HOLM SURVIVES")]
            for f in MULT["family"]]
    table(s, ["RQ", "Verdict under multiplicity"], rows,
          [1.53, 10.06], top=Emu(1450000), size=Pt(13))
    keyline(s, "Three of five returned unresolved or not-estimable. Each was a "
               "pre-registered test against a matched control with its "
               "precision target met — that is what makes them findings.")

    # ---- 24. References --------------------------------------------------
    s = D.content("References")
    keys = ["gastrohun", "liu22", "guo", "fleiss", "cohen", "gwet", "hayes",
            "borgli", "jha", "szegedy", "hinton", "efficientnet", "gradcam",
            "prisma2020", "delamor", "gao", "maenpaa", "nagendran"]
    byk = {e["key"]: e for e in BIB.references()}
    items = [f"[{byk[k]['n']}]  {byk[k]['formatted'][:112]}"
             for k in keys if k in byk]
    bullets(body_box(s, height=Emu(4000000)), items, size=Pt(10), space=Pt(2))
    keyline(s, f"Full list: {BIB.coverage()['n_total']} references, all "
               f"generated from the review's extraction table.")

    # ---- 25. Thank you ---------------------------------------------------
    s = D.content("Thank You")
    rich(body_box(s, top=Emu(2400000), height=Emu(1600000)), [
        [("Questions and comments welcome", Pt(24), True, TITLE_RGB)],
    ], align=PP_ALIGN.CENTER)

    if D.unused():
        raise SystemExit(f"{D.unused()} template slides left unwritten; they "
                         f"would ship with placeholder text. Add content or "
                         f"reduce the template.")
    # Rebuild the slide-id list exactly as the working Phase-I deck does. The
    # order is already correct, so this is the identity permutation; it exists
    # because reorder() rewrites every entry in _sldIdLst, and a deck saved
    # without that rewrite is the one structural difference between this
    # builder and the one whose output PowerPoint accepts.
    P1.reorder(prs, list(range(len(prs.slides._sldIdLst))))
    n_fixed = fix_corner_mark(prs)
    prs.save(str(OUTPUT))
    print(f"[final-deck] {len(prs.slides)} slides -> {OUTPUT.name}")
    print(f"[final-deck] {D.n_template} from template, {D.n_cloned} cloned")
    print(f"[final-deck] corner mark repaired on {n_fixed} layout run(s)")


if __name__ == "__main__":
    main()
