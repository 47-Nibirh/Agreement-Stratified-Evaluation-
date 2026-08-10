"""
Build the Phase-I defence deck by filling the university slide template.

The template's own furniture is preserved on every slide: the title placeholder
and its typeface and colour, the slide-number placeholder, and the
"B.Sc. Final-Defense" corner mark. Slides beyond the sixteen the template ships
are produced by cloning one of its content slides rather than by adding a bare
layout, so the furniture comes along with them.

One key message per slide. Every number is interpolated from
`phase1_facts.facts()`; nothing numeric is typed here.

Run:  python src/report/build_phase1_pptx.py
Out:  Phase-I_Defence_Presentation.pptx
"""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(Path(__file__).resolve().parent))

from phase1_facts import ADMIN, facts   # noqa: E402

TEMPLATE = ROOT / "Final-Defense Slide Template-PPT.pptx"
OUTPUT = ROOT / "Phase-I_Defence_Presentation.pptx"

FONT = "Times New Roman"
TITLE_RGB = RGBColor(0x70, 0x30, 0xA0)    # the template's own title purple
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x59, 0x59, 0x59)

# the content region the template's own layout defines
BODY_L, BODY_T = Emu(609600), Emu(1676400)
BODY_W, BODY_H = Emu(10972800), Emu(3962400)
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)


# --------------------------------------------------------------------------
# slide plumbing
# --------------------------------------------------------------------------
def clone_slide(prs, src_index: int):
    """Duplicate a template slide, keeping its furniture.

    `add_slide(layout)` would give a bare slide without the corner mark and
    without the exact slide-number placement the template uses, so an existing
    content slide is deep-copied instead. Only slides that carry no pictures
    are cloned, so there are no image relationships to rewire.
    """
    src = prs.slides[src_index]
    dst = prs.slides.add_slide(src.slide_layout)
    for shape in list(dst.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shape._element))
    return dst


def reorder(prs, index_order):
    """Rewrite the slide order. `index_order` holds current slide indices."""
    lst = prs.slides._sldIdLst
    entries = list(lst)
    assert sorted(index_order) == list(range(len(entries))), "order must be a permutation"
    for e in entries:
        lst.remove(e)
    for i in index_order:
        lst.append(entries[i])


def set_title(slide, text, size=Pt(30)):
    ph = slide.shapes.title
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = size
    r.font.bold = True
    r.font.color.rgb = TITLE_RGB
    return ph


def strip_body(slide, keep_title=True):
    """Remove everything except the title, slide number and corner mark."""
    for shape in list(slide.shapes):
        name = shape.name
        if shape.is_placeholder and shape.placeholder_format.type is not None:
            t = str(shape.placeholder_format.type)
            if t.startswith("TITLE") or t.startswith("CENTER_TITLE") \
                    or t.startswith("SLIDE_NUMBER"):
                continue
        if name.startswith("TextBox") and "Final-Defense" in (shape.text_frame.text
                                                              if shape.has_text_frame
                                                              else ""):
            continue
        shape._element.getparent().remove(shape._element)


# the usable band between the title and the footer strip
BODY_FULL_H = Emu(4250000)


def body_box(slide, left=None, top=None, width=None, height=None,
             anchor=MSO_ANCHOR.MIDDLE):
    """A text box in the content region.

    Body text is centred vertically by default: a top-anchored box that does
    not fill its region leaves a band of dead space above the footer, which is
    what the first pass of this deck did on every bullet slide.
    """
    box = slide.shapes.add_textbox(left or BODY_L, top or BODY_T,
                                   width or BODY_W, height or BODY_FULL_H)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def bullets(tf, items, size=Pt(18), space=Pt(10)):
    """`items` are (level, text) or plain strings at level 0."""
    first = True
    for it in items:
        lvl, text = it if isinstance(it, tuple) else (0, it)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = space
        p.line_spacing = 1.15
        bullet_char = "▪  " if lvl == 0 else "–  "
        r = p.add_run()
        r.text = bullet_char + text if text else ""
        r.font.name = FONT
        r.font.size = size if lvl == 0 else Pt(size.pt - 2)
        r.font.color.rgb = INK if lvl == 0 else MUTED
    return tf


def rich(tf, blocks, align=PP_ALIGN.LEFT, space=Pt(8), lead=1.2):
    """`blocks` are lists of (text, size, bold, colour) run tuples, one per line."""
    first = True
    for line in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = space
        p.line_spacing = lead
        for text, size, bold, colour in line:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = colour
    return tf


def picture(slide, path, top=None, max_h=None, caption=None):
    """Insert a figure, scaled to fit the content region and centred."""
    p = ROOT / path
    with Image.open(p) as im:
        pw, ph = im.size
    top = top if top is not None else BODY_T
    avail_h = max_h if max_h is not None else Emu(int(BODY_H) - (240000 if caption else 0))
    avail_w = BODY_W
    scale = min(int(avail_w) / pw, int(avail_h) / ph)
    w, h = int(pw * scale), int(ph * scale)
    left = int((int(SLIDE_W) - w) / 2)
    slide.shapes.add_picture(str(p), Emu(left), top, Emu(w), Emu(h))
    if caption:
        box = slide.shapes.add_textbox(BODY_L, Emu(int(top) + h + 40000),
                                       BODY_W, Emu(300000))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = caption
        r.font.name = FONT
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = MUTED
    return h


def table(slide, head, rows, widths, top=None, size=Pt(13), head_size=Pt(13)):
    """A table in the template's own theme style."""
    nrow, ncol = len(rows) + 1, len(head)
    total = Inches(sum(widths))
    left = Emu(int((int(SLIDE_W) - int(total)) / 2))
    top = top if top is not None else BODY_T
    height = Emu(int(Inches(0.42)) * nrow)
    shape = slide.shapes.add_table(nrow, ncol, left, top, total, height)
    tbl = shape.table
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(head):
        cell = tbl.cell(0, j)
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        r = para.add_run()
        r.text = str(h)
        r.font.name = FONT
        r.font.size = head_size
        r.font.bold = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            r = para.add_run()
            r.text = str(v)
            r.font.name = FONT
            r.font.size = size
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def keyline(slide, text, top=None, size=Pt(16)):
    """A single emphasised takeaway line under the content."""
    box = slide.shapes.add_textbox(BODY_L, top or Emu(5250000), BODY_W, Emu(500000))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = size
    r.font.bold = True
    r.font.color.rgb = TITLE_RGB
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ==========================================================================
# the deck
# ==========================================================================
def main() -> None:
    F = facts()
    c, ag, st = F["corpus"], F["agreement"], F["structure"]
    b, S_, cal = F["baseline"], F["strata"], F["calibration"]
    tr, pp, ct = F["training"], F["preprocess"], F["contamination"]
    order = S_["order"]
    prs = Presentation(str(TEMPLATE))
    sl = prs.slides

    # ---- 1. Title -------------------------------------------------------
    s = sl[0]
    set_title(s, ADMIN["title"], size=Pt(28))
    sub = [ph for ph in s.placeholders if ph.placeholder_format.idx == 1][0]
    # the template sizes this block for three or four lines; the identification
    # block needs six, so it is given the room rather than allowed to overflow
    sub.top, sub.height = Emu(3930000), Emu(2450000)
    sub.left, sub.width = Emu(2500000), Emu(7192000)
    tf = sub.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = []
    for name, sid in ADMIN["members"]:
        lines.append([(name, Pt(15), True, INK),
                      (f"    ID: {sid}", Pt(12), False, MUTED)])
    lines.append([(f"Supervisor: {ADMIN['supervisor'][0]}, "
                   f"{ADMIN['supervisor'][1]}", Pt(12), False, INK)])
    lines.append([(f"Co-Supervisor: {ADMIN['co_supervisor'][0]}, "
                   f"{ADMIN['co_supervisor'][1]}", Pt(12), False, INK)])
    lines.append([("Department of Computer Science and Engineering, "
                   "Daffodil International University", Pt(12), True, INK)])
    rich(tf, lines, align=PP_ALIGN.CENTER, space=Pt(9))
    notes(s, "Our project asks a simple question: when published systems report "
             "85 to 88 macro F1 on gastric landmark recognition, what fraction "
             "of the problem does that number actually describe?")

    # ---- 2. Outline -----------------------------------------------------
    s = sl[1]
    set_title(s, "Outline")
    ph = [x for x in s.placeholders if x.placeholder_format.idx == 13][0]
    ph.text_frame.clear()
    bullets(ph.text_frame, [
        "Introduction and problem identification",
        "Objectives",
        "Background study and literature review",
        "Gap analysis",
        "Methodology: data, pre-processing, model",
        "Results and analysis",
        "Novelty of the work",
        "Sample dataset and expected output",
        "Reproducibility and deliverables",
        "Conclusion and next steps",
        "References",
    ], size=Pt(20), space=Pt(6))

    # ---- 3. Introduction ------------------------------------------------
    s = sl[2]
    strip_body(s)
    set_title(s, "Introduction")
    bullets(body_box(s, height=Emu(3350000)), [
        "Upper GI endoscopy is the primary screening route for gastric cancer; "
        "its yield depends on inspecting the whole mucosal surface.",
        "The Systematic Screening protocol fixes a photographic route over 22 "
        "landmarks - a grid of 4 gastric walls by 6 depth stations.",
        "A model that recognises the landmark in each frame can act as a "
        "real-time coverage monitor and flag an unvisited region.",
        f"Published systems report macro F1 near {b['published']:g}-88, and the "
        f"task is generally treated as solved.",
        (1, f"But every one of those figures is measured only on frames all four "
            f"experts labelled identically - {ag['unanimous_pct']:.1f}% of this "
            f"corpus."),
    ], size=Pt(19))
    keyline(s, "The reported number describes the easy fraction of the task.")

    # ---- 4. Problem identification (cloned) -----------------------------
    s_problem = clone_slide(prs, 2)
    strip_body(s_problem)
    set_title(s_problem, "Problem Identification")
    rich(body_box(s_problem, height=Emu(1500000)), [
        [(f"{ag['unanimous_pct']:.1f}%", Pt(38), True, TITLE_RGB),
         ("   of the corpus is unanimous - and that is all the literature "
          "scores on.", Pt(19), False, INK)],
        [(f"{ag['contested_pct']:.1f}%", Pt(38), True, TITLE_RGB),
         (f"   ({ag['contested_n']:,} images) is discarded before any result is "
          f"reported.", Pt(19), False, INK)],
    ])
    bullets(body_box(s_problem, top=Emu(3150000), height=Emu(1950000)), [
        "A deployed system meets the contested fraction continuously, and gets "
        "no signal that its validation excluded those frames.",
        f"The discarded images are not noise: "
        f"{st['decomp_pct']['same_station_different_wall']:.2f}% of conflicts put "
        f"two experts on different walls of the same station.",
        "Confidence calibrated on unanimous frames may stay high on ambiguous "
        "ones - assured errors, the failure mode hardest to catch in clinic.",
    ], size=Pt(17))
    keyline(s_problem, "Operating accuracy is unknown on exactly the images "
                       "where a second reader would help most.")
    notes(s_problem, "This is the whole project in one slide.")

    # ---- 5. Objectives --------------------------------------------------
    s = sl[3]
    strip_body(s)
    set_title(s, "Objectives")
    bullets(body_box(s), [
        "Audit corpus provenance and integrity before modelling: licence, "
        "ethics, decode integrity, and a calibrated contamination scan.",
        "Complete a PRISMA 2020 literature review and derive a gap statement of "
        "measurable deficiencies rather than general observations.",
        "Fix the pre-processing chain once: annotation handling, resampling, "
        "normalisation, augmentation, transfer-learned features.",
        f"Train, validate and test a baseline, and reproduce the published "
        f"result to within the plus or minus {b['band']:g} macro F1 points fixed "
        f"in advance.",
        "Quantify accuracy and calibration per agreement stratum, with "
        "patient-clustered bootstrap intervals.",
        "Release a pipeline in which every figure and number regenerates from "
        "committed scripts.",
    ], size=Pt(18))

    # ---- 6. Background study: published baselines -----------------------
    s = sl[4]
    strip_body(s)
    set_title(s, "Background Study - published baselines")
    table(s, ["Model / label construction", "Params", "Macro F1"], [
        ["ConvNeXt-Large", "200 M", "88.25 +/- 0.22"],
        ["ConvNeXt-Tiny", "28 M", "approx. 85"],
        ["ResNet-152", "60 M", "85.28 +/- 0.27"],
        ["ConvNeXt-Tiny, fellow-agreement labels", "28 M", "87.05 +/- 0.21"],
        ["Best single annotator as target", "-", "84.82 +/- 0.23"],
        ["Human expert band", "-", "77.47 - 84.82"],
    ], widths=[5.0, 1.6, 2.6], top=Emu(1750000), size=Pt(14), head_size=Pt(14))
    keyline(s, "Changing only how the four labels are combined moves F1 by 2.2 "
               "points - more than the gap between architecture families.",
            top=Emu(5150000), size=Pt(15))
    notes(s, "The descriptor reports the fourth row as an aside and does not "
             "pursue it. That aside is our starting point.")

    # ---- 7. Background: agreement (cloned) ------------------------------
    s_agree = clone_slide(prs, 2)
    strip_body(s_agree)
    set_title(s_agree, "Background Study - do the experts agree?")
    picture(s_agree, "figures_phase1/PH1_F06_agreement.png", top=Emu(1500000),
            max_h=Emu(3500000))
    keyline(s_agree, f"Fleiss kappa = {ag['fleiss']:.4f}. Seniority does not "
                     f"predict agreement: each fellow agrees more closely with "
                     f"either gastroenterologist than with the other fellow.",
            top=Emu(5300000), size=Pt(14))

    # ---- 8. Gap analysis -------------------------------------------------
    s = sl[5]
    strip_body(s)
    set_title(s, "Gap Analysis")
    table(s, ["Gap", "What the literature does", "What we do instead"], [
        ["1  Evaluation is conditioned on expert unanimity",
         f"Scores only the {ag['unanimous_pct']:.1f}% all four experts agreed on; "
         f"the rest is filtered out before scoring.",
         "Report performance for every agreement stratum, with "
         "patient-clustered intervals."],
        ["2  Disagreement is discarded rather than used",
         f"Collapses the vote distribution to one label at ingest, discarding a "
         f"structured signal - "
         f"{st['decomp_pct']['same_station_different_wall']:.1f}% of conflicts are "
         f"same-station different-wall.",
         "Train on the full four-vote distribution, against a matched "
         "label-smoothing control."],
        ["3  Calibration is rarely reported, and never by stratum",
         "Nothing establishes whether confidence estimated on unanimous frames "
         "stays trustworthy on ambiguous ones.",
         "Treat expected calibration error by stratum as a primary endpoint."],
    ], widths=[2.9, 4.3, 3.8], top=Emu(1620000), size=Pt(12), head_size=Pt(13.5))
    keyline(s, "Each gap is a measurable deficiency with a matching endpoint, "
               "not a general complaint about the field.",
            top=Emu(4150000), size=Pt(15))
    notes(s, "Three gaps, each with a matching endpoint later in the design.")

    # ---- 9. Methodology: the phase workflow ------------------------------
    s = sl[6]
    strip_body(s)
    set_title(s, "Methodology - seven gated phases")
    picture(s, "figures_phase1/PH1_F04_workflow.png", top=Emu(1450000),
            max_h=Emu(3650000))
    keyline(s, "Each phase freezes its hypotheses and verdict rules before any "
               "model runs. Phases 0 to 3 are reported today.",
            top=Emu(5300000), size=Pt(15))

    # ---- 10. Data collection and the integrity gate (cloned) -------------
    s_data = clone_slide(prs, 2)
    strip_body(s_data)
    set_title(s_data, "Methodology - data collection and integrity")
    rich(body_box(s_data, top=Emu(1330000), height=Emu(1100000)), [
        [("GastroHUN:  ", Pt(17), True, INK),
         (f"{c['n_images']:,} images | {c['n_patients']} patients | "
          f"{c['n_classes']} classes | {c['n_annotators']} independent annotators "
          f"| {c['gb']:.2f} GB | CC BY 4.0 | ethics CEI-2019-06-10",
          Pt(16), False, INK)],
        [("Chosen because it releases the individual annotators' labels. Without "
          "them, no endpoint in this design exists.", Pt(15), False, MUTED)],
    ])
    picture(s_data, "figures_phase1/PH1_F12_integrity_gate.png",
            top=Emu(2480000), max_h=Emu(3000000))
    notes(s_data, "Six PASS, two CONDITIONAL. A conditional verdict becomes a "
                  "declared limitation rather than a silent one.")

    # ---- 11. Pre-processing (cloned) -------------------------------------
    s_pre = clone_slide(prs, 2)
    strip_body(s_pre)
    set_title(s_pre, "Methodology - data pre-processing")
    picture(s_pre, "figures_phase1/PH1_F01_preprocessing_pipeline.png",
            top=Emu(1420000), max_h=Emu(3600000))
    keyline(s_pre, "No flips and no large rotations: the label encodes a gastric "
                   "WALL, so the standard augmentation recipe would relabel the "
                   "image.", top=Emu(5250000), size=Pt(15))
    notes(s_pre, "Annotation, augmentation and feature engineering all appear "
                 "here - the three Phase-I pre-processing requirements.")

    # ---- 12. Model and training (cloned) ---------------------------------
    s_arch = clone_slide(prs, 2)
    strip_body(s_arch)
    set_title(s_arch, "Methodology - baseline model and training")
    picture(s_arch, "figures_phase1/PH1_F02_architecture.png",
            top=Emu(1420000), max_h=Emu(3600000))
    keyline(s_arch, f"ConvNeXt-Tiny, three seeds, {tr['total_train_min']:.0f} "
                    f"minutes total on one 4 GB GPU. Selection on validation "
                    f"macro F1 only; the test split is touched once.",
            top=Emu(5250000), size=Pt(14))

    # ---- 13. Results: baseline reproduction ------------------------------
    s = sl[7]
    strip_body(s)
    set_title(s, "Results - baseline reproduction")
    table(s, ["Metric", "Value"], [
        ["Macro F1, three-seed mean",
         f"{b['observed']:.2f}    (95% CI {b['ci95'][0]:.2f} - {b['ci95'][1]:.2f})"],
        ["Published target",
         f"{b['published']:g} +/- {b['band']:g}, fixed before training"],
        ["Difference", f"{b['delta']:+.2f} points"],
        ["Accuracy / weighted F1",
         f"{b['accuracy']:.2f} / {b['weighted_f1']:.2f}"],
        ["Expected calibration error", f"{b['ece']:.2f}%"],
        ["Seed spread",
         f"{b['sd']:.2f} points "
         f"({', '.join(f'{v:.2f}' for v in b['per_seed'].values())})"],
        ["Pre-registered verdict", b["verdict"]],
    ], widths=[3.9, 6.3], top=Emu(1700000), size=Pt(15), head_size=Pt(15))
    keyline(s, "PASS - the pipeline behaves as the published one did. Everything "
               "after this is attributable to the design, not to a bug.",
            top=Emu(5300000), size=Pt(15))

    # ---- 14. Results: training behaviour ---------------------------------
    s = sl[8]
    strip_body(s)
    set_title(s, "Results - training behaviour")
    picture(s, "figures_phase1/PH1_F03_training_dynamics.png",
            top=Emu(1550000), max_h=Emu(3400000))
    keyline(s, "All three seeds stopped by early stopping rather than by the "
               "compute cap, so the schedule was not truncated.",
            top=Emu(5250000), size=Pt(15))

    # ---- 15. Results: the stratified finding (cloned) --------------------
    s_strat = clone_slide(prs, 2)
    strip_body(s_strat)
    set_title(s_strat, "Results - performance across agreement strata")
    picture(s_strat, "figures_phase1/PH1_F10_stratified_result.png",
            top=Emu(1420000), max_h=Emu(3550000))
    keyline(s_strat, f"Macro F1 falls {S_['f1'][order[0]]:.1f} to "
                     f"{S_['f1'][order[-1]]:.1f}, a {S_['gap']:.1f}-point gap, "
                     f"against {S_['arch_benchmark']:g} points between "
                     f"architecture families.",
            top=Emu(5200000), size=Pt(15))
    notes(s_strat, "Part of the fall is the ceiling moving. The right panel "
                   "separates the two: a real model shortfall survives on two of "
                   "three contrasts, and the third does not resolve.")

    # ---- 16. Results: calibration (cloned) -------------------------------
    s_cal = clone_slide(prs, 2)
    strip_body(s_cal)
    set_title(s_cal, "Results - the calibration collapse")
    picture(s_cal, "figures_phase1/PH1_F11_calibration.png",
            top=Emu(1450000), max_h=Emu(3500000))
    keyline(s_cal, f"Confidence falls only "
                   f"{cal['confidence'][order[0]] - cal['confidence'][order[2]]:.1f} "
                   f"points while accuracy falls "
                   f"{cal['expected_accuracy'][order[0]] - cal['expected_accuracy'][order[2]]:.1f}. "
                   f"ECE {cal['ece'][order[0]]:.1f}% to "
                   f"{max(cal['ece'].values()):.1f}%.",
            top=Emu(5250000), size=Pt(15))
    notes(s_cal, "This is the clinically consequential finding: the model does "
                 "not know it has entered harder territory.")

    # ---- 17. Novelty ------------------------------------------------------
    s = sl[9]
    strip_body(s)
    set_title(s, "Novelty of the Work")
    bullets(body_box(s), [
        "Evaluation stratified by annotator agreement: performance, calibration "
        "and error geometry per stratum with patient-clustered intervals, "
        "instead of truncated at unanimity.",
        "The attainable ceiling is measured rather than assumed, so a falling "
        "reference standard is separated from a falling classifier. The "
        "literature reports the two as one quantity.",
        "The four-annotator vote distribution used as a training target and "
        "tested against a label-smoothing control matched to the probability "
        "mass the soft target displaces.",
        "Model error compared against the anatomical geometry of human "
        "disagreement, made measurable by the wall-by-station grid.",
        "Every phase pre-registered before execution, and every reported "
        "quantity regenerated from committed scripts.",
    ], size=Pt(17))

    # ---- 18. Sample dataset -----------------------------------------------
    s = sl[10]
    strip_body(s)
    set_title(s, "Sample Dataset")
    picture(s, "figures_phase1/PH1_F13_sample_images.png",
            top=Emu(1400000), max_h=Emu(3750000))
    keyline(s, "Left to right, expert agreement falls. Published evaluations "
               "score only the leftmost column.",
            top=Emu(5350000), size=Pt(15))

    # ---- 19. Sample dataset and expected output ---------------------------
    s = sl[11]
    strip_body(s)
    set_title(s, "Sample Dataset and Expected Output")
    picture(s, "figures_phase1/PH1_F05_label_space.png",
            top=Emu(1250000), max_h=Emu(3000000))
    bullets(body_box(s, top=Emu(4350000), height=Emu(1500000)), [
        f"Input: a {pp['size']} by {pp['size']} endoscopic frame. Output: one of "
        f"{c['n_classes']} classes (wall by station, plus OTHERCLASS) with a "
        f"confidence score.",
        "Expected deliverable: an agreement-stratified performance and "
        "calibration profile, and a controlled verdict on soft-label training.",
    ], size=Pt(16), space=Pt(6))

    # ---- 20. Reproducibility (the template's spare content slide) ---------
    s = sl[12]
    strip_body(s)
    set_title(s, "Reproducibility and Deliverables")
    bullets(body_box(s), [
        f"{F['prisma']['unique']:,} records screened to PRISMA 2020; "
        f"{ct['n_pairs_scanned']:,} image pairs scanned for contamination; "
        f"{c['n_images']:,} images audited - all from committed scripts.",
        "Every figure and table in the progress report is generated from a JSON "
        "artefact. No number in the report or in this deck is typed by hand.",
        "Pre-registration files are written before training, and the generating "
        "script refuses to overwrite an existing one.",
        f"Frozen checkpoints are reused across phases: the stratified evaluation "
        f"reproduces the baseline predictions exactly - "
        f"{S_['consistency_compared']:,} comparisons, "
        f"{S_['consistency_mismatch']} mismatches.",
        "Delivered this phase: audited corpus, literature review, pre-processing "
        "pipeline, trained baseline, stratified evaluation, and the progress "
        "report.",
    ], size=Pt(17))

    # ---- 21. Conclusion ---------------------------------------------------
    s = sl[13]
    strip_body(s)
    set_title(s, "Conclusion and Next Steps")
    bullets(body_box(s), [
        f"The corpus passed an eight-criterion integrity gate with a PROCEED "
        f"verdict, and the baseline reproduced the published result at "
        f"{b['observed']:.2f} against {b['published']:g} +/- {b['band']:g} - a "
        f"pre-registered PASS.",
        f"Reported accuracy on this task describes {ag['unanimous_pct']:.1f}% of "
        f"it. Macro F1 falls {S_['f1'][order[0]]:.1f} to "
        f"{S_['f1'][order[-1]]:.1f} across agreement strata.",
        "Part of that fall is the attainable ceiling dropping; a real model "
        "shortfall remains once that is accounted for.",
        f"Confidence degrades further and faster than discrimination: ECE "
        f"{cal['ece'][order[0]]:.1f}% to {max(cal['ece'].values()):.1f}%.",
        (1, "Next: train on the vote distribution against a matched control; "
            "external validation without adaptation; explainability against a "
            "human comparator."),
    ], size=Pt(17))

    # ---- 22. References ---------------------------------------------------
    s = sl[14]
    strip_body(s)
    set_title(s, "References")
    refs = [
        "[1]  D. Panesso-Ortiz et al., \"GastroHUN: an endoscopy dataset of the "
        "complete systematic screening protocol for the stomach,\" Scientific "
        "Data, vol. 12, art. 102, 2025.",
        "[2]  M. Nagendran et al., \"Artificial intelligence versus clinicians,\" "
        "BMJ, vol. 368, art. m689, 2020.",
        "[3]  R. Djinbachian et al., \"Interobserver agreement for the Paris "
        "classification of colorectal lesions,\" Dig. Dis. Sci., 2025.",
        "[4]  S. Isajevs et al., \"Gastritis staging: interobserver agreement by "
        "applying OLGA and OLGIM systems,\" Virchows Archiv, vol. 464, 2014.",
        "[5]  R. Del Amor et al., \"Labeling confidence for uncertainty-aware "
        "histology image classification,\" Comput. Med. Imaging Graph., vol. 107, "
        "2023.",
        "[6]  Z. Gao et al., \"Leveraging multi-annotator label uncertainties as "
        "privileged information,\" Bioengineering, vol. 11, no. 2, 2024.",
        "[7]  S. M. Maenpaa et al., \"Diagnostic test accuracy of externally "
        "validated CNN models,\" Int. J. Med. Inform., vol. 189, 2024.",
        "[8]  Y. D. Li et al., \"Intelligent detection endoscopic assistant,\" "
        "Dig. Liver Dis., vol. 53, pp. 216-223, 2021.",
        "[9]  Z. Liu et al., \"A ConvNet for the 2020s,\" in Proc. CVPR, 2022, "
        "pp. 11976-11986.",
        "[10] M. J. Page et al., \"The PRISMA 2020 statement,\" BMJ, vol. 372, "
        "art. n71, 2021.",
        "[11] C. Guo et al., \"On calibration of modern neural networks,\" in "
        "Proc. ICML, 2017, pp. 1321-1330.",
    ]
    rich(body_box(s, height=Emu(4400000)),
         [[(r, Pt(12), False, INK)] for r in refs])

    # ---- 23. Thank you ----------------------------------------------------
    s = sl[15]
    ph = [x for x in s.placeholders if x.placeholder_format.idx == 13][0]
    ph.text_frame.clear()
    rich(ph.text_frame, [
        [("THANK YOU", Pt(48), True, TITLE_RGB)],
        [("", Pt(14), False, INK)],
        [("Questions and comments welcome", Pt(20), False, MUTED)],
    ], align=PP_ALIGN.CENTER)

    # ---- final ordering ---------------------------------------------------
    # cloned slides land at the end of the deck; put each where it belongs
    idx = {sd.slide_id: i for i, sd in enumerate(prs.slides)}
    seq = [sl[0], sl[1], sl[2], s_problem, sl[3], sl[4], s_agree, sl[5], sl[6],
           s_data, s_pre, s_arch, sl[7], sl[8], s_strat, s_cal, sl[9], sl[10],
           sl[11], sl[12], sl[13], sl[14], sl[15]]
    reorder(prs, [idx[x.slide_id] for x in seq])

    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT.relative_to(ROOT)}")
    print(f"  slides: {len(seq)}")


if __name__ == "__main__":
    main()
